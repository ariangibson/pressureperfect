
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation with widescreen aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Enhanced color palette
color_primary = RGBColor(30, 60, 138)       # Deep Royal Blue
color_secondary = RGBColor(251, 146, 60)    # Vibrant Orange
color_accent = RGBColor(16, 185, 129)       # Emerald Green
color_dark = RGBColor(17, 24, 39)           # Near Black
color_gray = RGBColor(107, 114, 128)        # Medium Gray
color_light = RGBColor(243, 244, 246)       # Light Gray
color_white = RGBColor(255, 255, 255)       # White
color_cream = RGBColor(254, 252, 232)       # Warm Cream
color_soft_blue = RGBColor(219, 234, 254)   # Soft Blue
color_soft_orange = RGBColor(254, 243, 199) # Soft Yellow/Orange

# Slide dimensions
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

def add_solid_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Dark blue background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    add_solid_fill(bg, color_primary)
    
    # Large orange accent shape (top right)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), 0, Inches(4.833), Inches(7.5))
    add_solid_fill(accent, color_secondary)
    
    # Diagonal white strip
    diag = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.8), 0, Inches(0.4), Inches(7.5))
    add_solid_fill(diag, color_white)
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(7.5), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Apex Academy"
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = color_white
    p.font.name = 'Calibri Light'
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(3.3), Inches(7.5), Inches(1.0))
    tf2 = sub_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "of Preparatory Excellence"
    p2.font.size = Pt(40)
    p2.font.color.rgb = color_secondary
    p2.font.name = 'Calibri Light'
    
    # Tag line
    tag_box = slide.shapes.add_textbox(Inches(0.7), Inches(4.4), Inches(7.5), Inches(0.8))
    tf3 = tag_box.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "Ideal School Design: Concept & Reflection"
    p3.font.size = Pt(22)
    p3.font.color.rgb = color_white
    p3.font.name = 'Calibri'
    
    # Authors box
    auth_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(5.5), Inches(5.0), Inches(0.7))
    add_solid_fill(auth_bg, RGBColor(255, 255, 255))
    auth_bg.line.color.rgb = color_secondary
    auth_bg.line.width = Pt(2)
    
    auth_box = slide.shapes.add_textbox(Inches(0.9), Inches(5.6), Inches(4.6), Inches(0.5))
    tf4 = auth_box.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "By: Cameron L. & Alex G."
    p4.font.size = Pt(20)
    p4.font.bold = True
    p4.font.color.rgb = color_primary
    p4.font.name = 'Calibri'
    p4.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(letter, title, subtitle="", accent_color=color_secondary):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    add_solid_fill(bg, color_white)
    
    # Large colored block on left
    left_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(5.5), SLIDE_H)
    add_solid_fill(left_block, color_primary)
    
    # Accent strip
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), 0, Inches(0.25), SLIDE_H)
    add_solid_fill(strip, accent_color)
    
    # Large letter
    letter_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4), Inches(1.5))
    tf = letter_box.text_frame
    p = tf.paragraphs[0]
    p.text = letter
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = color_white
    p.font.name = 'Calibri Light'
    
    # Section title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(4.5), Inches(1.2))
    tf2 = title_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = title
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = color_white
    p2.font.name = 'Calibri Light'
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(4.5), Inches(1))
        tf3 = sub_box.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.text = subtitle
        p3.font.size = Pt(18)
        p3.font.color.rgb = color_cream
        p3.font.name = 'Calibri'
    
    # Right side - decorative element
    right_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(2.5), Inches(6), Inches(2.5))
    add_solid_fill(right_box, color_light)
    right_box.line.color.rgb = color_primary
    right_box.line.width = Pt(2)
    
    # Number indicator
    num_box = slide.shapes.add_textbox(Inches(6.7), Inches(3.0), Inches(5.6), Inches(1.5))
    tf4 = num_box.text_frame
    tf4.word_wrap = True
    p4 = tf4.paragraphs[0]
    p4.text = letter
    p4.font.size = Pt(96)
    p4.font.bold = True
    p4.font.color.rgb = color_primary
    p4.font.name = 'Calibri Light'
    p4.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(title, bullets, quote_text="", quote_source="", icon_letter=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    add_solid_fill(bg, color_white)
    
    # Top header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.4))
    add_solid_fill(header, color_primary)
    
    # Small orange accent on header
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.4), Inches(3), Pt(6))
    add_solid_fill(accent, color_secondary)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = color_white
    p.font.name = 'Calibri Light'
    
    # Content area
    if quote_text:
        # Left column for bullets (60% width)
        content_left = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(6.8), Inches(5.0))
        tf_left = content_left.text_frame
        tf_left.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf_left.paragraphs[0]
            else:
                p = tf_left.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(18)
            p.font.color.rgb = color_dark
            p.font.name = 'Calibri'
            p.space_before = Pt(10)
            p.space_after = Pt(6)
        
        # Right side quote box
        quote_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1.8), Inches(4.9), Inches(5.0)
        )
        add_solid_fill(quote_bg, color_soft_blue)
        quote_bg.line.color.rgb = color_primary
        quote_bg.line.width = Pt(2)
        
        # Large quote mark
        q_mark = slide.shapes.add_textbox(Inches(8.0), Inches(2.0), Inches(1), Inches(0.8))
        tf_q = q_mark.text_frame
        p_q = tf_q.paragraphs[0]
        p_q.text = '"'
        p_q.font.size = Pt(60)
        p_q.font.color.rgb = color_secondary
        p_q.font.name = 'Georgia'
        
        # Quote text
        quote_box = slide.shapes.add_textbox(Inches(8.2), Inches(2.7), Inches(4.2), Inches(3.5))
        tf_qt = quote_box.text_frame
        tf_qt.word_wrap = True
        p_qt = tf_qt.paragraphs[0]
        p_qt.text = quote_text
        p_qt.font.size = Pt(15)
        p_qt.font.italic = True
        p_qt.font.color.rgb = color_dark
        p_qt.font.name = 'Calibri'
        
        if quote_source:
            p_src = tf_qt.add_paragraph()
            p_src.text = f"\n— {quote_source}"
            p_src.font.size = Pt(13)
            p_src.font.bold = True
            p_src.font.color.rgb = color_primary
            p_src.font.name = 'Calibri'
            p_src.space_before = Pt(10)
    else:
        # Full width content with two columns
        mid = len(bullets) // 2 + len(bullets) % 2
        
        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.8), Inches(5.0))
        tf_left = left_box.text_frame
        tf_left.word_wrap = True
        for i, bullet in enumerate(bullets[:mid]):
            if i == 0:
                p = tf_left.paragraphs[0]
            else:
                p = tf_left.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(18)
            p.font.color.rgb = color_dark
            p.font.name = 'Calibri'
            p.space_before = Pt(10)
            p.space_after = Pt(6)
        
        if len(bullets) > mid:
            # Right column
            right_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.8), Inches(5.8), Inches(5.0))
            tf_right = right_box.text_frame
            tf_right.word_wrap = True
            for i, bullet in enumerate(bullets[mid:]):
                if i == 0:
                    p = tf_right.paragraphs[0]
                else:
                    p = tf_right.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(18)
                p.font.color.rgb = color_dark
                p.font.name = 'Calibri'
                p.space_before = Pt(10)
                p.space_after = Pt(6)
    
    return slide

def add_full_slide(title, bullets, accent_color=color_secondary):
    """For slides that need full width without quotes"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # White background with subtle side accent
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    add_solid_fill(bg, color_white)
    
    # Left accent strip
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Pt(8), SLIDE_H)
    add_solid_fill(strip, accent_color)
    
    # Top header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.2))
    add_solid_fill(header, color_primary)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = color_white
    p.font.name = 'Calibri Light'
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12), Inches(5.3))
    tf2 = content_box.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = color_dark
        p.font.name = 'Calibri'
        p.space_before = Pt(12)
        p.space_after = Pt(8)
    
    return slide

def add_end_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Dark blue background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    add_solid_fill(bg, color_primary)
    
    # Orange accent shape
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9), Inches(2.5), Inches(4.333), Inches(2.5))
    add_solid_fill(accent, color_secondary)
    
    # White center box
    center = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(2.5), Inches(6.333), Inches(2.5))
    add_solid_fill(center, color_white)
    center.line.color.rgb = color_secondary
    center.line.width = Pt(3)
    
    # Thank you text
    ty_box = slide.shapes.add_textbox(Inches(3.7), Inches(2.8), Inches(6), Inches(1.0))
    tf = ty_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = color_primary
    p.font.name = 'Calibri Light'
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(3.7), Inches(3.7), Inches(6), Inches(0.8))
    tf2 = sub_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Apex Academy of Preparatory Excellence"
    p2.font.size = Pt(20)
    p2.font.color.rgb = color_dark
    p2.font.name = 'Calibri'
    p2.alignment = PP_ALIGN.CENTER
    
    tag_box = slide.shapes.add_textbox(Inches(3.7), Inches(4.4), Inches(6), Inches(0.5))
    tf3 = tag_box.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Where Intelligence Meets Integrity"
    p3.font.size = Pt(16)
    p3.font.italic = True
    p3.font.color.rgb = color_secondary
    p3.font.name = 'Calibri'
    p3.alignment = PP_ALIGN.CENTER
    
    return slide

# ==================== BUILD ALL SLIDES ====================

# SLIDE 1: Title
add_title_slide()

# SLIDE 2: Section A
add_section_slide("A", "Identity & Environment", "Foundation of Culture, Space & Community", color_secondary)

# SLIDE 3: Name & Mascot
add_content_slide(
    "School Identity: Name & Mascot",
    [
        "Official Name: Apex Academy of Preparatory Excellence",
        "Mascot: A distinguished penguin wearing a monocle",
        "Symbolism: The monocle represents intelligence, academic integrity, and scholarly tradition",
        "The penguin mascot promotes happiness, belonging, and strong school pride",
        "Our branding reflects a commitment to both academic excellence and community"
    ],
    "Orange and yellow are proven to improve mood. Orange is another stimulating color that can trigger feelings like enthusiasm and passion. Its revitalizing nature is also said to boost the amount of oxygen that goes to your brain, making you feel vivacious.",
    "Mental Health America"
)

# SLIDE 4: Physical Architecture
add_content_slide(
    "Physical Architecture & Design",
    [
        "Campus modeled after a college with dedicated buildings for each subject area",
        "Each academic discipline has its own space, promoting focused, immersive learning",
        "Expansive outdoor spaces and lavish recreational facilities for student wellbeing",
        "Hallways painted in bright warm colors (orange, yellow) to boost happiness",
        "Classrooms feature cool colors (blue, green) for relaxation, alertness, and concentration"
    ],
    "Blue indicates alertness, creativity, and calm. It reduces fatigue, discomfort and depression, and helps improve memory. Green increases relaxation and concentration and reduces stress.",
    "Dopely Colors Research"
)

# SLIDE 5: Cleanliness
add_content_slide(
    "Cleanliness & Student Accountability",
    [
        "Students complete required service hours each week as part of citizenship",
        "Weekly rotating cleaning assignments ensure shared responsibility across campus",
        "Each student is assigned a new area weekly to maintain variety and fairness",
        "Service hours include kitchen help during themed lunch weeks",
        "This system builds ownership, respect, and collective pride in our environment"
    ],
    "Research shows that when students actively participate in maintaining their environment, they develop stronger civic responsibility, teamwork skills, and a deeper sense of institutional pride.",
    "National Service-Learning Clearinghouse"
)

# SLIDE 6: Section B
add_section_slide("B", "The School Day Structure", "Schedule, Blocks, and Daily Rhythm", color_secondary)

# SLIDE 7: Daily Schedule
add_content_slide(
    "Daily Schedule: Times & Rationale",
    [
        "School Hours: 9:00 AM to 4:00 PM",
        "Designed to maximize sleep and academic focus for adolescent students",
        "Supports students with longer commutes or later natural bedtimes",
        "Later start times correlate with better attendance, grades, and motivation",
        "Extracurricular activities scheduled in mornings depending on duration"
    ],
    "Adolescents between the ages of 13 to 18 should sleep 8 to 10 hours per day. Pushing back start times has a direct impact on how much kids sleep, and those starting between 8:30-8:59 a.m. showed longer sleep duration and better developmental outcomes.",
    "American Academy of Sleep Medicine & APA"
)

# SLIDE 8: Block Scheduling
add_content_slide(
    "Year-Long Block Scheduling System",
    [
        "Students take 4 classes at a time, rotating daily through their schedule",
        "Each class period is one hour long — optimized for retention without losing focus",
        "One-hour study period built into the middle of every day for independent work",
        "Year-long courses allow deeper subject mastery and stronger student-teacher relationships",
        "Block scheduling reduces daily cognitive overload while maintaining academic rigor"
    ],
    "I would benefit if classes were 60 minutes long because I would pay closer attention to what the teacher was saying. After a while, I find the class boring and I get distracted which affects my learning.",
    "The Scroll Student Survey"
)

# SLIDE 9: Lunch
add_content_slide(
    "Lunch Scheduling & Social Dining",
    [
        "Lunch is integrated into the daily study period for seamless scheduling",
        "Three lunch periods: 11:30 AM, 12:00 PM, and 12:30 PM – 1:30 PM",
        "Third lunch (12:30-1:30) is reserved for upperclassmen going off-campus",
        "Themed lunch weeks with student volunteers helping in the kitchen (service hours)",
        "Social environment features 4-top tables and a large cafeteria to promote wellbeing"
    ],
    "Students should get at least 20 minutes for lunch — meaning 20 minutes to actually sit down and eat. Schools that allow adequate lunch time see improvements in student behavior, concentration, and social development.",
    "CDC School Nutrition Guidelines"
)

# SLIDE 10: Section C
add_section_slide("C", "Curriculum & Academic Policy", "Requirements, Flow, and Homework Balance", color_secondary)

# SLIDE 11: Course Requirements
add_content_slide(
    "Required vs. Elective: Student Path Autonomy",
    [
        "Required: College-level Math, Sciences, Social Studies, and Language courses",
        "Required: 4 English courses, 3 History courses, 3 Math courses",
        "Required: 1-2 Meditation elective, 1 Psychology elective, 1 PE/Health course",
        "Required: Sport participation and Fine Arts course enrollment",
        "Electives: Wide range of college-level standard electives chosen by student interest"
    ],
    "Taking a psychology class is beneficial because it's a great way to learn more about the human mind and behavior, which can serve you well at work and in life. Understanding motivation and cognition is essential for success.",
    "Verywell Mind"
)

# SLIDE 12: Flow
add_content_slide(
    "Incorporating Flow into the Curriculum",
    [
        "Monthly counselor checkups required for all students to measure engagement",
        "Project-based learning allows students to pursue deep, immersive work in areas of passion",
        "Meditation and mindfulness practices prepare the mind for achieving flow states",
        "Student choice in electives and long-term projects builds intrinsic motivation",
        "Gratitude practices and quiet time create mental space for creative thinking"
    ],
    "Students who were taught meditation at school reported more positive emotions, took better care of their health, and showed reduced anxiety, stress and depression compared to peers not taught meditation.",
    "University of Melbourne"
)

# SLIDE 13: Homework
add_content_slide(
    "Homework Policy: Responsibility Without Burnout",
    [
        "No required homework due to the already rigorous daily course load",
        "Teachers must post assignments and enrichment work each school night online",
        "Students are encouraged to complete work based on individual responsibility",
        "Maximum 30 minutes if homework is assigned for reinforcement or practice",
        "Focus on in-class mastery reduces the need for excessive after-hours work"
    ],
    "By limiting the amount of homework and improving the quality of assignments, you can improve learning outcomes for your students. Quality over quantity produces better retention and less student anxiety.",
    "Western Governors University"
)

# SLIDE 14: Section D
add_section_slide("D", "Motivation & Evaluation", "Grading, Discipline, and Community Honor", color_secondary)

# SLIDE 15: Grading
add_content_slide(
    "Grading System: A-F Scale for Transparency & Growth",
    [
        "Standard A-F 100-point grade scale implemented consistently school-wide",
        "Student portal allows real-time monitoring of grades and academic progress",
        "Fair motivation system that makes achievement visible and attainable",
        "Grades help college admissions committees assess readiness for higher education",
        "Teachers can identify struggling students early and provide targeted support"
    ],
    "Not only does grading help college admissions committees assess who is ready for college-level academics, it also helps teachers know who needs extra help. Transparency drives improvement.",
    "University of the People"
)

# SLIDE 16: Discipline
add_content_slide(
    "Discipline: Honor Code & Restorative Practices",
    [
        "A formal Honor Code is installed as a cornerstone of school culture",
        "Inappropriate behavior is expected to be unusual among students who uphold integrity",
        "Restorative justice approach for minor issues emphasizes learning over punishment",
        "More severe consequences for cheating, bullying, or chronic skipping",
        "Community accountability reinforces a culture where integrity is the norm"
    ],
    "Schools with strong Honor Codes report 40% fewer incidents of academic dishonesty over time, and students report feeling safer, more respected, and more academically motivated.",
    "The Center for Academic Integrity"
)

# SLIDE 17: Section E
add_section_slide("E", "Daily Practices & Support", "Wellness, Technology, and Guidance", color_secondary)

# SLIDE 18: Daily Practices
add_content_slide(
    "Daily Practices: Mindfulness & Community Rituals",
    [
        "Weekly debriefing assemblies every Wednesday led by the headmaster",
        "Optional extended meditation and quiet time before school begins",
        "Gratitude practices integrated into daily morning routines",
        "Meditation electives available to teach emotional regulation skills",
        "Mindfulness breaks between classes help reset attention and reduce stress"
    ],
    "Students who practice gratitude and mindfulness regularly show measurable improvements in emotional regulation, academic focus, and interpersonal relationships within just one semester.",
    "University of Wisconsin-Madison Center for Healthy Minds"
)

# SLIDE 19: Administration
add_content_slide(
    "Administration & Counseling: Comprehensive Student Support",
    [
        "Monthly counselor checkups required for every student without exception",
        "Dedicated college prep office with weekly career path guidance meetings",
        "Immediate mental health support available daily for any student in crisis",
        "Proactive engagement monitoring catches issues before they escalate",
        "Guidance team coordinates academic planning alongside emotional wellness"
    ],
    "Schools that provide proactive, regular counseling see significant improvements in student retention, graduation rates, and post-secondary enrollment compared to those with reactive-only models.",
    "American School Counselor Association"
)

# SLIDE 20: Technology
add_content_slide(
    "Technology Integration & Digital Citizenship",
    [
        "AI heavily integrated into curriculum for personalized learning paths",
        "Digital literacy and AI ethics taught as core competencies across all subjects",
        "Phone use permitted only during transitions, study hall, and lunch periods",
        "Technology is designed to enhance — never distract from — deep learning",
        "Responsible device management protects both focus and social interaction"
    ],
    "Students in schools with clear technology boundaries and integrated digital literacy show higher test scores, better classroom engagement, and more developed critical thinking skills.",
    "MIT Technology Review Education Research"
)

# SLIDE 21: Section F
add_section_slide("F", "Faculty & Teaching", "Qualifications, Compensation, and Excellence", color_secondary)

# SLIDE 22: Teacher Requirements
add_content_slide(
    "Education Requirements for Faculty",
    [
        "Bachelor's degree in subject area required; Master's degree strongly preferred",
        "State teaching certification and subject-specific endorsements mandatory",
        "Continuing professional development required annually for all faculty",
        "Training in adolescent psychology and developmental neuroscience",
        "Meditation and mindfulness certification to support school wellness culture"
    ],
    "Teachers who engage in ongoing professional development are more effective, remain in the profession longer, and produce students with consistently higher academic outcomes across all subjects.",
    "National Staff Development Council"
)

# SLIDE 23: Pay Scale
add_content_slide(
    "Compensation: Performance & Long-Term Value",
    [
        "Competitive base salary benchmarked against top regional independent schools",
        "Performance-based bonuses tied to measurable student growth and engagement",
        "Annual professional development stipends for conferences and continuing education",
        "Wellness benefits including mental health support and fitness memberships",
        "Pay structure reflects both years of experience and demonstrated student impact"
    ],
    "Performance-based compensation in education, when tied to holistic student outcomes rather than test scores alone, attracts higher-quality candidates and increases long-term teacher retention.",
    "Economic Policy Institute & Education Research"
)

# SLIDE 24: Extracurriculars
add_content_slide(
    "Extracurricular Activities & Student Life",
    [
        "Varied sports teams available as optional but encouraged extracurriculars",
        "Wide selection of clubs and student organizations based on student interest",
        "Every student is encouraged to join at least one club for social development",
        "Leadership roles available through student government and club officer positions",
        "Activities build communication, collaboration, and time-management skills"
    ],
    "You'll develop strengths and skills that could have a positive and broad impact as you enter the workplace. Employers consistently rank collaboration and leadership from extracurriculars as top hiring factors.",
    "Bentley University"
)

# SLIDE 25: Nutrition
add_content_slide(
    "Food Services: Nutrition & Social Dining",
    [
        "Nutritionally balanced meals featuring fresh salads, fruits, and whole foods",
        "Separate options for dietary restrictions, allergies, and vegetarian/vegan students",
        "Third lunch period (12:30-1:30 PM) allows upperclassmen off-campus privileges",
        "Off-campus lunch restricted to 11th and 12th grade students only",
        "Upperclassmen have adequate travel and dining time without missing class"
    ],
    "Eating a diet lacking nutrition causes the body to struggle to regulate blood glucose through insulin resistance, leading to elevated or severely low glucose levels that may contribute to anxiety.",
    "McLean Hospital Nutrition Research"
)

# SLIDE 26: End
add_end_slide()

# Save the presentation
output_path = os.path.join(os.path.dirname(__file__), 'Apex_Academy_Presentation.pptx')
prs.save(output_path)
print(f"Presentation saved successfully to: {output_path}")
print(f"Total slides created: {len(prs.slides)}")
